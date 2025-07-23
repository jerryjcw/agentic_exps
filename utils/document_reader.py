#!/usr/bin/env python3
"""
Document Reader - Universal document reading utility.

This module provides a comprehensive document reader that supports multiple file formats:
- PDF files (.pdf)
- Word documents (.docx)
- PowerPoint presentations (.pptx)
- Text files (.txt)
- Markdown files (.md)
- CSV files (.csv)
- Excel files (.xlsx)
"""

import os
import logging
import re
import unicodedata
from pathlib import Path
from typing import Dict, List, Union, Optional, Tuple

# Document format dependencies
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd

# GPU detection and Docling imports
try:
    import torch
    from docling.document_converter import DocumentConverter
    DOCLING_AVAILABLE = True
except ImportError:
    torch = None
    DocumentConverter = None
    DOCLING_AVAILABLE = False

# Configure logging
logger = logging.getLogger(__name__)


class DocumentReader:
    """
    Universal document reader supporting multiple file formats.
    
    Supported formats:
    - PDF (.pdf) - using Docling (GPU) or PyPDF2 (CPU fallback)
    - Word (.docx) - using Docling (GPU) or python-docx (CPU fallback)
    - PowerPoint (.pptx) - using Docling (GPU) or python-pptx (CPU fallback)
    - Text (.txt) - native Python
    - Markdown (.md) - native Python
    - CSV (.csv) - using pandas
    - Excel (.xlsx) - using pandas
    """
    
    def __init__(self, use_gpu: Optional[bool] = None):
        """
        Initialize the DocumentReader with format handlers.
        
        Args:
            use_gpu: If True, force GPU usage. If False, force CPU. If None, auto-detect.
        """
        self.use_docling = self._should_use_docling(use_gpu)
        
        if self.use_docling:
            self._init_docling()
            logger.info("Initialized DocumentReader with Docling (GPU acceleration)")
        else:
            logger.info("Initialized DocumentReader with traditional libraries (CPU fallback)")
        
        self.supported_formats = {
            '.pdf': self._read_pdf,
            '.docx': self._read_docx,
            '.pptx': self._read_pptx,
            '.txt': self._read_text,
            '.md': self._read_text,
            '.csv': self._read_csv,
            '.xlsx': self._read_xlsx
        }
    
    def read_document(self, file_path: Union[str, Path], as_markdown: bool = False) -> str:
        """
        Read a single document and return its content as string.
        
        Args:
            file_path: Path to the document file
            as_markdown: If True, format output as markdown preserving structure
            
        Returns:
            str: Document content as text or markdown
            
        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file format is not supported
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}. "
                           f"Supported formats: {list(self.supported_formats.keys())}")
        
        try:
            content = self.supported_formats[file_extension](file_path, as_markdown)
            logger.info(f"Successfully read {file_extension} file: {file_path.name}")
            return content
        except Exception as e:
            logger.error(f"Error reading {file_path}: {e}")
            raise
    
    def read_multiple_documents(self, file_paths: List[Union[str, Path]], as_markdown: bool = False) -> Dict[str, Dict[str, str]]:
        """
        Read multiple documents and return a mapping of results.
        
        Args:
            file_paths: List of paths to document files
            as_markdown: If True, format output as markdown preserving structure
            
        Returns:
            Dict[str, Dict[str, str]]: Dictionary mapping file names to their metadata:
                {
                    "filename.pdf": {
                        "file_type": "pdf",
                        "content": "document content...",
                        "status": "success"
                    },
                    "failed_file.doc": {
                        "file_type": "doc", 
                        "content": "",
                        "status": "error",
                        "error": "error message"
                    }
                }
        """
        results = {}
        
        for file_path in file_paths:
            file_path = Path(file_path)
            file_name = file_path.name
            file_type = file_path.suffix.lower().lstrip('.')
            
            try:
                content = self.read_document(file_path, as_markdown)
                results[file_name] = {
                    "file_type": file_type,
                    "content": content,
                    "status": "success"
                }
            except Exception as e:
                results[file_name] = {
                    "file_type": file_type,
                    "content": "",
                    "status": "error",
                    "error": str(e)
                }
                logger.error(f"Failed to read {file_name}: {e}")
        
        return results
    
    def _should_use_docling(self, use_gpu: Optional[bool] = None) -> bool:
        """
        Determine whether to use Docling based on GPU availability and configuration.
        
        Args:
            use_gpu: User preference for GPU usage
            
        Returns:
            bool: True if Docling should be used
        """
        if not DOCLING_AVAILABLE:
            logger.info("Docling not available, using traditional libraries")
            return False
        
        if use_gpu is False:
            logger.info("GPU usage explicitly disabled, using traditional libraries")
            return False
        
        # Check if GPU is available
        gpu_available = torch and torch.cuda.is_available()
        
        if use_gpu is True and not gpu_available:
            logger.warning("GPU usage requested but no GPU available, falling back to traditional libraries")
            return False
        
        if gpu_available:
            logger.info(f"GPU detected: {torch.cuda.get_device_name(0) if torch.cuda.is_available() else 'Unknown'}")
            return True
        
        # Check for Apple Silicon MPS
        if torch and hasattr(torch.backends, 'mps') and torch.backends.mps.is_available():
            logger.info("Apple Silicon MPS detected, using Docling")
            return True
        
        logger.info("No GPU acceleration available, using traditional libraries")
        return False
    
    def _init_docling(self):
        """
        Initialize Docling with optimal GPU configuration.
        """
        try:
            # Determine best device and set GPU memory management
            if torch.cuda.is_available():
                device_name = torch.cuda.get_device_name(0)
                logger.info(f"Using CUDA device: {device_name}")
                
                # Set GPU memory management
                torch.cuda.set_per_process_memory_fraction(0.8)
            elif hasattr(torch.backends, 'mps') and torch.backends.mps.is_available():
                logger.info("Using Apple Silicon MPS")
            else:
                logger.info("Using CPU acceleration")
            
            # Initialize converter - using default initialization for broad compatibility
            self.docling_converter = DocumentConverter()
            
        except Exception as e:
            logger.warning(f"Failed to initialize Docling: {e}, falling back to CPU")
            self.use_docling = False
    
    def get_supported_formats(self) -> List[str]:
        """
        Get list of supported file formats.
        
        Returns:
            List[str]: List of supported file extensions
        """
        return list(self.supported_formats.keys())
    
    def is_using_gpu(self) -> bool:
        """
        Check if the reader is using GPU acceleration.
        
        Returns:
            bool: True if using GPU acceleration (Docling)
        """
        return self.use_docling
    
    def is_supported(self, file_path: Union[str, Path]) -> bool:
        """
        Check if a file format is supported.
        
        Args:
            file_path: Path to check
            
        Returns:
            bool: True if format is supported
        """
        file_extension = Path(file_path).suffix.lower()
        return file_extension in self.supported_formats
    
    def _detect_list_item(self, line: str) -> Optional[str]:
        """
        Detect if a line is a list item and return the formatted markdown version.
        
        This method uses a systematic approach to detect various bullet point and numbering patterns:
        - Unicode bullet characters (General_Category = Symbol other)
        - Common ASCII bullet patterns 
        - Numbered lists with various delimiters
        - Lettered lists
        - Roman numerals
        - Indented patterns
        
        Args:
            line: The line to analyze
            
        Returns:
            str: Formatted markdown list item, or None if not a list item
        """
        line = line.strip()
        if not line:
            return None
        
        # Pattern 1: Unicode bullet characters (comprehensive)
        # Check if the first character is a bullet-like symbol
        first_char = line[0]
        if len(line) > 1:
            # Unicode categories that typically contain bullet characters
            char_category = unicodedata.category(first_char)
            if char_category in ['So', 'Sm', 'Sc']:  # Symbol other, Symbol math, Symbol currency
                # Additional check: common bullet Unicode ranges
                char_code = ord(first_char)
                bullet_ranges = [
                    (0x2022, 0x2043),  # General Punctuation bullets
                    (0x25A0, 0x25FF),  # Geometric Shapes 
                    (0x2190, 0x21FF),  # Arrows
                    (0x2600, 0x26FF),  # Miscellaneous Symbols
                    (0x2700, 0x27BF),  # Dingbats
                ]
                if any(start <= char_code <= end for start, end in bullet_ranges):
                    return f"- {line[1:].strip()}"
        
        # Pattern 2: ASCII bullet patterns with regex
        ascii_bullet_pattern = r'^[-*+•◦▪▫‣⁃] '
        if re.match(ascii_bullet_pattern, line):
            return f"- {line[1:].strip()}"
        
        # Pattern 3: Numbered lists (comprehensive)
        # Arabic numerals: 1. 1) 1: 1- 1/ 1\
        numbered_pattern = r'^(\d+)[.)\-:/\\]\s+'
        match = re.match(numbered_pattern, line)
        if match:
            return f"1. {line[len(match.group(0)):].strip()}"
        
        # Pattern 4: Lettered lists: a. A. a) A) a: A:
        letter_pattern = r'^([a-zA-Z])[.)\-:]\s+'
        match = re.match(letter_pattern, line)
        if match:
            return f"1. {line[len(match.group(0)):].strip()}"
        
        # Pattern 5: Roman numerals: i. I. ii. II. iii. III. iv. IV. v. V.
        roman_pattern = r'^([ivxlcdmIVXLCDM]+)[.)\-:]\s+'
        match = re.match(roman_pattern, line)
        if match:
            # Verify it's actually a roman numeral
            roman_chars = set('ivxlcdmIVXLCDM')
            if all(c in roman_chars for c in match.group(1)):
                return f"1. {line[len(match.group(0)):].strip()}"
        
        # Pattern 6: Indented bullet points (no explicit bullet character)
        # Lines that start with 2+ spaces followed by text could be sub-bullets
        indent_pattern = r'^(\s{2,})([^\s].*)'
        match = re.match(indent_pattern, line)
        if match:
            indent_level = len(match.group(1)) // 2  # 2 spaces = 1 indent level
            prefix = "  " * max(0, indent_level - 1)  # Markdown nested list format
            return f"{prefix}- {match.group(2)}"
        
        # Pattern 7: Special characters that commonly serve as bullets
        special_bullets = ['→', '➤', '▶', '►', '‣', '⁃', '⦿', '⦾', '○', '●', '◯', '◉']
        if any(line.startswith(bullet + ' ') for bullet in special_bullets):
            for bullet in special_bullets:
                if line.startswith(bullet + ' '):
                    return f"- {line[len(bullet):].strip()}"
        
        return None
    
    # Format-specific readers
    
    def _read_pdf(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read PDF file using Docling (GPU) or PyPDF2 (CPU fallback)."""
        if self.use_docling:
            return self._read_pdf_docling(file_path, as_markdown)
        else:
            return self._read_pdf_pypdf2(file_path, as_markdown)
    
    def _read_pdf_docling(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read PDF file using Docling with GPU acceleration."""
        try:
            result = self.docling_converter.convert(file_path)
            
            if as_markdown:
                return result.document.export_to_markdown()
            else:
                # Convert to plain text with page separators for compatibility
                markdown_content = result.document.export_to_markdown()
                return self._convert_docling_markdown_to_plain_with_pages(markdown_content, result.document.num_pages)
        except Exception as e:
            # Clear GPU cache and try CPU fallback
            if torch and torch.cuda.is_available():
                torch.cuda.empty_cache()
            logger.warning(f"Docling PDF processing failed: {e}, falling back to PyPDF2")
            return self._read_pdf_pypdf2(file_path, as_markdown)
    
    def _read_pdf_pypdf2(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read PDF file using PyPDF2 (original implementation)."""
        content = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            if as_markdown:
                                # Format as markdown with page headers
                                formatted_text = self._format_pdf_text_as_markdown(text, page_num + 1)
                                content.append(formatted_text)
                            else:
                                content.append(f"--- Page {page_num + 1} ---\n{text}")
                    except Exception as e:
                        logger.warning(f"Could not extract text from page {page_num + 1}: {e}")
                        if as_markdown:
                            content.append(f"## Page {page_num + 1}\n\n*[Text extraction failed]*")
                        else:
                            content.append(f"--- Page {page_num + 1} ---\n[Text extraction failed]")
        except Exception as e:
            raise Exception(f"Error reading PDF file: {e}")
        
        return "\n\n".join(content) if content else "[No readable text found in PDF]"
    
    def _format_pdf_text_as_markdown(self, text: str, page_num: int) -> str:
        """Format PDF text as markdown with advanced structure detection."""
        lines = text.split('\n')
        formatted_lines = [f"## Page {page_num}"]
        
        for line in lines:
            original_line = line
            line = line.strip()
            if not line:
                continue
                
            # Detect potential headings (all caps, short lines)
            if len(line) < 80 and line.isupper() and len(line) > 3:
                formatted_lines.append(f"\n### {line.title()}")
            else:
                # Use systematic list detection
                list_item = self._detect_list_item(original_line)
                if list_item:
                    formatted_lines.append(list_item)
                else:
                    # Regular paragraph text
                    formatted_lines.append(line)
        
        return '\n'.join(formatted_lines)
    
    def _format_docx_paragraph_as_markdown(self, paragraph) -> str:
        """Format a Word paragraph as markdown based on its style."""
        text = paragraph.text.strip()
        if not text:
            return ""
            
        # Check paragraph style for headings
        style_name = paragraph.style.name.lower() if paragraph.style else ""
        
        # Handle headings
        if 'heading 1' in style_name or 'title' in style_name:
            return f"# {text}"
        elif 'heading 2' in style_name:
            return f"## {text}"
        elif 'heading 3' in style_name:
            return f"### {text}"
        elif 'heading 4' in style_name:
            return f"#### {text}"
        elif 'heading 5' in style_name:
            return f"##### {text}"
        elif 'heading 6' in style_name:
            return f"###### {text}"
        
        # Handle lists using systematic detection
        list_item = self._detect_list_item(text)
        if list_item:
            return list_item
        
        # Check for bold/italic formatting in runs
        formatted_text = ""
        for run in paragraph.runs:
            run_text = run.text
            if run.bold and run.italic:
                formatted_text += f"***{run_text}***"
            elif run.bold:
                formatted_text += f"**{run_text}**"
            elif run.italic:
                formatted_text += f"*{run_text}*"
            else:
                formatted_text += run_text
        
        return formatted_text if formatted_text.strip() else text
    
    def _format_table_as_markdown(self, table) -> str:
        """Format a table as markdown."""
        if not table.rows:
            return ""
            
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            
            # Add header separator after first row
            if i == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                rows.append(separator)
        
        return "\n".join(rows)
    
    def _format_pptx_text_as_markdown(self, text: str) -> str:
        """Format PowerPoint text as markdown with advanced list detection."""
        lines = text.split('\n')
        formatted_lines = []
        
        for line in lines:
            original_line = line
            line = line.strip()
            if not line:
                continue
                
            # Detect potential slide titles (usually first non-empty line or all caps)
            if len(line) < 100 and (line.isupper() or len(formatted_lines) == 0):
                formatted_lines.append(f"### {line}")
            else:
                # Use systematic list detection
                list_item = self._detect_list_item(original_line)
                if list_item:
                    formatted_lines.append(list_item)
                else:
                    # Regular text
                    formatted_lines.append(line)
        
        return '\n'.join(formatted_lines)
    
    def _format_pptx_table_as_markdown(self, table) -> str:
        """Format a PowerPoint table as markdown."""
        if not table.rows:
            return ""
            
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            
            # Add header separator after first row
            if i == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                rows.append(separator)
        
        return "\n".join(rows)
    
    def _read_docx(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read Word document using Docling (GPU) or python-docx (CPU fallback)."""
        if self.use_docling:
            return self._read_docx_docling(file_path, as_markdown)
        else:
            return self._read_docx_python_docx(file_path, as_markdown)
    
    def _read_docx_docling(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read Word document using Docling with GPU acceleration."""
        try:
            result = self.docling_converter.convert(file_path)
            
            if as_markdown:
                return result.document.export_to_markdown()
            else:
                # Convert to plain text for compatibility
                return result.document.export_to_text()
        except Exception as e:
            # Clear GPU cache and try CPU fallback
            if torch and torch.cuda.is_available():
                torch.cuda.empty_cache()
            logger.warning(f"Docling DOCX processing failed: {e}, falling back to python-docx")
            return self._read_docx_python_docx(file_path, as_markdown)
    
    def _read_docx_python_docx(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read Word document using python-docx (original implementation)."""
        try:
            doc = Document(file_path)
            content = []
            
            # Extract paragraphs with style information for markdown
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    if as_markdown:
                        formatted_text = self._format_docx_paragraph_as_markdown(paragraph)
                        content.append(formatted_text)
                    else:
                        content.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                if as_markdown:
                    markdown_table = self._format_table_as_markdown(table)
                    content.append(markdown_table)
                else:
                    table_content = []
                    for row in table.rows:
                        row_content = []
                        for cell in row.cells:
                            row_content.append(cell.text.strip())
                        table_content.append(" | ".join(row_content))
                    if table_content:
                        content.append("--- Table ---\n" + "\n".join(table_content))
            
            return "\n\n".join(content) if content else "[No readable text found in document]"
            
        except Exception as e:
            raise Exception(f"Error reading Word document: {e}")
    
    def _read_pptx(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read PowerPoint presentation using Docling (GPU) or python-pptx (CPU fallback)."""
        if self.use_docling:
            return self._read_pptx_docling(file_path, as_markdown)
        else:
            return self._read_pptx_python_pptx(file_path, as_markdown)
    
    def _read_pptx_docling(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read PowerPoint presentation using Docling with GPU acceleration."""
        try:
            result = self.docling_converter.convert(file_path)
            
            if as_markdown:
                return result.document.export_to_markdown()
            else:
                # Convert to plain text with slide separators for compatibility
                markdown_content = result.document.export_to_markdown()
                return self._convert_docling_pptx_markdown_to_plain(markdown_content)
        except Exception as e:
            # Clear GPU cache and try CPU fallback
            if torch and torch.cuda.is_available():
                torch.cuda.empty_cache()
            logger.warning(f"Docling PPTX processing failed: {e}, falling back to python-pptx")
            return self._read_pptx_python_pptx(file_path, as_markdown)
    
    def _read_pptx_python_pptx(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read PowerPoint presentation using python-pptx (original implementation)."""
        try:
            prs = Presentation(file_path)
            content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                if as_markdown:
                    slide_content = [f"## Slide {slide_num}"]
                else:
                    slide_content = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if as_markdown:
                            # Format slide text with basic markdown
                            formatted_text = self._format_pptx_text_as_markdown(shape.text)
                            slide_content.append(formatted_text)
                        else:
                            slide_content.append(shape.text)
                    elif shape.has_table:
                        # Extract table content
                        if as_markdown:
                            markdown_table = self._format_pptx_table_as_markdown(shape.table)
                            slide_content.append(markdown_table)
                        else:
                            table_content = []
                            for row in shape.table.rows:
                                row_content = []
                                for cell in row.cells:
                                    row_content.append(cell.text.strip())
                                table_content.append(" | ".join(row_content))
                            if table_content:
                                slide_content.append("Table:\n" + "\n".join(table_content))
                
                if len(slide_content) > 1:  # More than just the slide header
                    content.append("\n".join(slide_content))
            
            return "\n\n".join(content) if content else "[No readable text found in presentation]"
            
        except Exception as e:
            raise Exception(f"Error reading PowerPoint presentation: {e}")
    
    def _read_text(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read text file (.txt, .md)."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        except UnicodeDecodeError:
            # Try with different encodings
            encodings = ['latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                        logger.warning(f"File read with {encoding} encoding instead of UTF-8")
                        break
                except UnicodeDecodeError:
                    continue
            else:
                raise Exception(f"Could not decode text file with any supported encoding")
        except Exception as e:
            raise Exception(f"Error reading text file: {e}")
        
        # Apply markdown formatting to .txt files (but not .md files which are already markdown)
        if as_markdown and file_path.suffix.lower() == '.txt':
            return self._format_text_as_markdown(content)
        
        return content
    
    def _format_text_as_markdown(self, text: str) -> str:
        """Format plain text as markdown with list detection."""
        lines = text.split('\n')
        formatted_lines = []
        
        for line in lines:
            if not line.strip():
                formatted_lines.append(line)  # Preserve empty lines
                continue
            
            # Use systematic list detection
            list_item = self._detect_list_item(line)
            if list_item:
                formatted_lines.append(list_item)
            else:
                # Regular text line
                formatted_lines.append(line.rstrip())
        
        return '\n'.join(formatted_lines)
    
    def _read_csv(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read CSV file using pandas."""
        try:
            # Try to read CSV with various parameters
            try:
                df = pd.read_csv(file_path)
            except UnicodeDecodeError:
                df = pd.read_csv(file_path, encoding='latin-1')
            except pd.errors.EmptyDataError:
                return "[Empty CSV file]"
            
            # Convert to string representation
            content = []
            
            if as_markdown:
                content.append(f"# CSV File: {file_path.name}")
                content.append(f"**Shape:** {df.shape[0]} rows, {df.shape[1]} columns")
                content.append(f"**Columns:** {', '.join(df.columns.tolist())}")
                content.append("")
                # Convert DataFrame to markdown table
                markdown_table = self._dataframe_to_markdown(df, max_rows=100)
                content.append(markdown_table)
                
                if len(df) > 100:
                    content.append(f"\n*... and {len(df) - 100} more rows*")
            else:
                content.append(f"CSV File: {file_path.name}")
                content.append(f"Shape: {df.shape[0]} rows, {df.shape[1]} columns")
                content.append(f"Columns: {', '.join(df.columns.tolist())}")
                content.append("\n--- Data ---")
                content.append(df.to_string(index=False, max_rows=100))
                
                if len(df) > 100:
                    content.append(f"\n... and {len(df) - 100} more rows")
            
            return "\n".join(content)
            
        except Exception as e:
            raise Exception(f"Error reading CSV file: {e}")
    
    def _read_xlsx(self, file_path: Path, as_markdown: bool = False) -> str:
        """Read Excel file using pandas."""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content = []
            if as_markdown:
                content.append(f"# Excel File: {file_path.name}")
                content.append(f"**Sheets:** {', '.join(excel_file.sheet_names)}")
            else:
                content.append(f"Excel File: {file_path.name}")
                content.append(f"Sheets: {', '.join(excel_file.sheet_names)}")
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    if not df.empty:
                        if as_markdown:
                            content.append(f"\n## Sheet: {sheet_name}")
                            content.append(f"**Shape:** {df.shape[0]} rows, {df.shape[1]} columns")
                            content.append(f"**Columns:** {', '.join(df.columns.astype(str).tolist())}")
                            content.append("")
                            markdown_table = self._dataframe_to_markdown(df, max_rows=50)
                            content.append(markdown_table)
                            
                            if len(df) > 50:
                                content.append(f"\n*... and {len(df) - 50} more rows*")
                        else:
                            content.append(f"\n--- Sheet: {sheet_name} ---")
                            content.append(f"Shape: {df.shape[0]} rows, {df.shape[1]} columns")
                            content.append(f"Columns: {', '.join(df.columns.astype(str).tolist())}")
                            content.append(df.to_string(index=False, max_rows=50))
                            
                            if len(df) > 50:
                                content.append(f"... and {len(df) - 50} more rows")
                    else:
                        if as_markdown:
                            content.append(f"\n## Sheet: {sheet_name}")
                            content.append("*[Empty sheet]*")
                        else:
                            content.append(f"\n--- Sheet: {sheet_name} ---")
                            content.append("[Empty sheet]")
                except Exception as sheet_error:
                    if as_markdown:
                        content.append(f"\n## Sheet: {sheet_name}")
                        content.append(f"*[Error reading sheet: {sheet_error}]*")
                    else:
                        content.append(f"\n--- Sheet: {sheet_name} ---")
                        content.append(f"[Error reading sheet: {sheet_error}]")
            
            return "\n".join(content)
            
        except Exception as e:
            raise Exception(f"Error reading Excel file: {e}")
    
    def _convert_docling_markdown_to_plain_with_pages(self, markdown_content: str, num_pages: int) -> str:
        """Convert Docling markdown output to plain text with page separators for compatibility."""
        lines = markdown_content.split('\n')
        plain_lines = []
        current_page = 1
        
        for line in lines:
            # Remove markdown formatting for plain text compatibility
            line = re.sub(r'^#+\s+', '', line)  # Remove headers
            line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)  # Remove bold
            line = re.sub(r'\*(.*?)\*', r'\1', line)  # Remove italic
            line = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', line)  # Remove links
            
            # Add page separators similar to PyPDF2 format
            if line.strip().startswith('<!-- image -->') or (current_page <= num_pages and len(plain_lines) > 0 and len(plain_lines) % 50 == 0):
                if current_page <= num_pages:
                    plain_lines.append(f"--- Page {current_page} ---")
                    current_page += 1
            
            if line.strip():
                plain_lines.append(line)
        
        return '\n'.join(plain_lines)
    
    def _convert_docling_pptx_markdown_to_plain(self, markdown_content: str) -> str:
        """Convert Docling PPTX markdown output to plain text with slide separators for compatibility."""
        lines = markdown_content.split('\n')
        plain_lines = []
        slide_num = 1
        
        for line in lines:
            # Detect slide boundaries and add separators
            if line.strip().startswith('##') and 'slide' not in line.lower():
                plain_lines.append(f"--- Slide {slide_num} ---")
                slide_num += 1
            
            # Remove markdown formatting for plain text compatibility
            line = re.sub(r'^#+\s+', '', line)  # Remove headers
            line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)  # Remove bold
            line = re.sub(r'\*(.*?)\*', r'\1', line)  # Remove italic
            line = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', line)  # Remove links
            
            if line.strip():
                plain_lines.append(line)
        
        return '\n'.join(plain_lines)
    
    def _dataframe_to_markdown(self, df, max_rows: int = None) -> str:
        """Convert pandas DataFrame to markdown table format."""
        if df.empty:
            return "*[Empty table]*"
        
        # Limit rows if specified
        display_df = df.head(max_rows) if max_rows else df
        
        # Convert all columns to string to handle various data types
        display_df = display_df.astype(str)
        
        # Create header row
        headers = list(display_df.columns)
        header_row = "| " + " | ".join(headers) + " |"
        
        # Create separator row
        separator_row = "| " + " | ".join(["---"] * len(headers)) + " |"
        
        # Create data rows
        data_rows = []
        for _, row in display_df.iterrows():
            # Clean cell values (remove newlines, limit length)
            cells = [str(cell).replace('\n', ' ').replace('|', '\\|')[:100] for cell in row]
            data_row = "| " + " | ".join(cells) + " |"
            data_rows.append(data_row)
        
        # Combine all rows
        return "\n".join([header_row, separator_row] + data_rows)
    
    def save_markdown(self, markdown_content: str, output_path: Union[str, Path]) -> None:
        """
        Save markdown content to a file.
        
        Args:
            markdown_content: The markdown-formatted content to save
            output_path: Path where the markdown file should be saved
            
        Raises:
            IOError: If the file cannot be written
            ValueError: If markdown_content is empty or None
        """
        if not markdown_content:
            raise ValueError("Markdown content cannot be empty or None")
        
        output_path = Path(output_path)
        
        # Ensure the output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Ensure the file has .md extension
        if output_path.suffix.lower() != '.md':
            output_path = output_path.with_suffix('.md')
        
        try:
            with open(output_path, 'w', encoding='utf-8') as file:
                file.write(markdown_content)
            logger.info(f"Successfully saved markdown content to: {output_path}")
        except Exception as e:
            logger.error(f"Error saving markdown file {output_path}: {e}")
            raise IOError(f"Failed to save markdown file: {e}")
    
    def read_and_save_as_markdown(self, input_path: Union[str, Path], output_path: Union[str, Path]) -> str:
        """
        Read a document and save it as markdown in one operation.
        
        Args:
            input_path: Path to the input document
            output_path: Path where the markdown file should be saved
            
        Returns:
            str: The markdown content that was saved
            
        Raises:
            FileNotFoundError: If input file doesn't exist
            ValueError: If file format is not supported
            IOError: If the output file cannot be written
        """
        # Read the document as markdown
        markdown_content = self.read_document(input_path, as_markdown=True)
        
        # Save the markdown content
        self.save_markdown(markdown_content, output_path)
        
        return markdown_content


# Convenience functions for direct usage
def read_document(file_path: Union[str, Path], as_markdown: bool = False, use_gpu: Optional[bool] = None) -> str:
    """
    Convenience function to read a single document.
    
    Args:
        file_path: Path to the document
        as_markdown: If True, format output as markdown preserving structure
        use_gpu: If True, force GPU usage. If False, force CPU. If None, auto-detect.
        
    Returns:
        str: Document content as text or markdown
    """
    reader = DocumentReader(use_gpu=use_gpu)
    return reader.read_document(file_path, as_markdown)


def read_multiple_documents(file_paths: List[Union[str, Path]], as_markdown: bool = False, use_gpu: Optional[bool] = None) -> Dict[str, Dict[str, str]]:
    """
    Convenience function to read multiple documents.
    
    Args:
        file_paths: List of file paths
        as_markdown: If True, format output as markdown preserving structure
        use_gpu: If True, force GPU usage. If False, force CPU. If None, auto-detect.
        
    Returns:
        Dict: Results mapping
    """
    reader = DocumentReader(use_gpu=use_gpu)
    return reader.read_multiple_documents(file_paths, as_markdown)


def get_supported_formats() -> List[str]:
    """
    Get list of supported file formats.
    
    Returns:
        List[str]: Supported file extensions
    """
    reader = DocumentReader()
    return reader.get_supported_formats()


def read_documents(file_paths, use_gpu: Optional[bool] = None):
    """Convenience function to read PDF documents."""
    reader = DocumentReader(use_gpu=use_gpu)
    content = reader.read_multiple_documents(file_paths, as_markdown=True)
    return content


# Example usage
if __name__ == "__main__":
    # Example usage
    reader = DocumentReader()
    print("Supported formats:", reader.get_supported_formats())
    docs = ["tests/resources/test.pdf", "tests/resources/test2.pdf", "tests/resources/test.docx"]
    for doc in docs:
        print(f'The document is {doc}')
        if reader.is_supported(doc):
            print(f"Reading {doc}...")
            content = reader.read_document(doc, as_markdown=True)
            print(f"Content of {doc}:\n{content[:500]}...")
        else:
            print(f"Skipping unsupported file format: {doc}")